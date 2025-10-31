"""Document parsing utilities for PDF and Word documents."""

import logging
import os
from pathlib import Path
from typing import Dict, List, Optional, Any, Union
import io

logger = logging.getLogger(__name__)


class DocumentParser:
    """Handles parsing of various document formats."""
    
    def __init__(self):
        """Initialize document parser."""
        self.supported_formats = {
            '.pdf': self._parse_pdf,
            '.docx': self._parse_docx,
            '.doc': self._parse_doc,
            '.txt': self._parse_text,
            '.md': self._parse_text,
            '.rtf': self._parse_rtf,
            '.pptx': self._parse_pptx,
            '.ppt': self._parse_ppt,
            # Media files - will extract text/transcripts where possible
            '.mp4': self._parse_media,
            '.mov': self._parse_media,
            '.m4a': self._parse_media,
            '.mp3': self._parse_media,
            '.wav': self._parse_media,
            '.avi': self._parse_media
        }
    
    def parse_document(self, file_path: Union[str, Path]) -> Dict[str, Any]:
        """Parse document and extract text content.
        
        Args:
            file_path: Path to document file
            
        Returns:
            Dictionary containing parsed content and metadata
        """
        path = Path(file_path)
        
        if not path.exists():
            raise FileNotFoundError(f"Document not found: {path}")
        
        file_extension = path.suffix.lower()
        
        if file_extension not in self.supported_formats:
            raise ValueError(f"Unsupported file format: {file_extension}")
        
        logger.info(f"Parsing document: {path}")
        
        try:
            parse_func = self.supported_formats[file_extension]
            result = parse_func(path)
            
            # Add metadata
            result['file_path'] = str(path)
            result['file_extension'] = file_extension
            result['file_size'] = path.stat().st_size
            
            logger.info(f"Successfully parsed document: {path}")
            return result
            
        except Exception as e:
            logger.error(f"Error parsing document {path}: {e}")
            raise
    
    def _parse_pdf(self, file_path: Path) -> Dict[str, Any]:
        """Parse PDF document.
        
        Args:
            file_path: Path to PDF file
            
        Returns:
            Dictionary with extracted content
        """
        try:
            import PyPDF2
        except ImportError:
            raise ImportError(
                "PyPDF2 is required for PDF parsing. Install with: pip install PyPDF2"
            )
        
        content = ""
        metadata = {}
        pages = []
        
        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            
            # Extract metadata
            if pdf_reader.metadata:
                metadata = {
                    'title': pdf_reader.metadata.get('/Title', ''),
                    'author': pdf_reader.metadata.get('/Author', ''),
                    'subject': pdf_reader.metadata.get('/Subject', ''),
                    'creator': pdf_reader.metadata.get('/Creator', ''),
                    'producer': pdf_reader.metadata.get('/Producer', ''),
                    'creation_date': str(pdf_reader.metadata.get('/CreationDate', '')),
                    'modification_date': str(pdf_reader.metadata.get('/ModDate', ''))
                }
            
            # Extract text from all pages
            for page_num, page in enumerate(pdf_reader.pages):
                try:
                    page_text = page.extract_text()
                    if page_text.strip():
                        pages.append({
                            'page_number': page_num + 1,
                            'text': page_text,
                            'word_count': len(page_text.split())
                        })
                        content += page_text + "\n\n"
                except Exception as e:
                    logger.warning(f"Error extracting text from page {page_num + 1}: {e}")
                    pages.append({
                        'page_number': page_num + 1,
                        'text': '',
                        'word_count': 0,
                        'error': str(e)
                    })
        
        return {
            'content': content.strip(),
            'metadata': metadata,
            'pages': pages,
            'total_pages': len(pages),
            'total_words': len(content.split())
        }
    
    def _parse_docx(self, file_path: Path) -> Dict[str, Any]:
        """Parse DOCX document.
        
        Args:
            file_path: Path to DOCX file
            
        Returns:
            Dictionary with extracted content
        """
        try:
            from docx import Document
        except ImportError:
            raise ImportError(
                "python-docx is required for DOCX parsing. Install with: pip install python-docx"
            )
        
        doc = Document(file_path)
        content = ""
        paragraphs = []
        tables = []
        
        # Extract metadata
        metadata = {}
        if doc.core_properties:
            metadata = {
                'title': doc.core_properties.title or '',
                'author': doc.core_properties.author or '',
                'subject': doc.core_properties.subject or '',
                'keywords': doc.core_properties.keywords or '',
                'created': str(doc.core_properties.created) if doc.core_properties.created else '',
                'modified': str(doc.core_properties.modified) if doc.core_properties.modified else '',
                'last_modified_by': doc.core_properties.last_modified_by or ''
            }
        
        # Extract paragraphs
        for para_num, paragraph in enumerate(doc.paragraphs):
            para_text = paragraph.text.strip()
            if para_text:
                paragraphs.append({
                    'paragraph_number': para_num + 1,
                    'text': para_text,
                    'style': paragraph.style.name if paragraph.style else '',
                    'word_count': len(para_text.split())
                })
                content += para_text + "\n\n"
        
        # Extract tables
        for table_num, table in enumerate(doc.tables):
            table_data = []
            for row in table.rows:
                row_data = []
                for cell in row.cells:
                    cell_text = cell.text.strip()
                    row_data.append(cell_text)
                table_data.append(row_data)
            
            tables.append({
                'table_number': table_num + 1,
                'data': table_data,
                'rows': len(table_data),
                'columns': len(table_data[0]) if table_data else 0
            })
        
        return {
            'content': content.strip(),
            'metadata': metadata,
            'paragraphs': paragraphs,
            'tables': tables,
            'total_paragraphs': len(paragraphs),
            'total_tables': len(tables),
            'total_words': len(content.split())
        }
    
    def _parse_doc(self, file_path: Path) -> Dict[str, Any]:
        """Parse DOC document (legacy Word format).
        
        Args:
            file_path: Path to DOC file
            
        Returns:
            Dictionary with extracted content
        """
        try:
            import subprocess
            import tempfile
        except ImportError:
            raise ImportError("subprocess and tempfile modules are required")
        
        # Try to convert DOC to text using antiword (Linux/Mac) or catdoc
        content = ""
        metadata = {}
        
        try:
            # Try antiword first (common on Linux/Mac)
            result = subprocess.run(
                ['antiword', str(file_path)],
                capture_output=True,
                text=True,
                timeout=30
            )
            if result.returncode == 0:
                content = result.stdout
            else:
                # Try catdoc as fallback
                result = subprocess.run(
                    ['catdoc', str(file_path)],
                    capture_output=True,
                    text=True,
                    timeout=30
                )
                if result.returncode == 0:
                    content = result.stdout
                else:
                    raise Exception("Unable to parse DOC file - install antiword or catdoc")
                    
        except (subprocess.TimeoutExpired, FileNotFoundError, Exception) as e:
            logger.warning(f"Could not parse DOC file with external tools: {e}")
            # Fallback: try to read as binary and extract what we can
            try:
                with open(file_path, 'rb') as f:
                    raw_content = f.read()
                    # Simple text extraction (very basic)
                    content = raw_content.decode('utf-8', errors='ignore')
                    # Clean up the content
                    import re
                    content = re.sub(r'[^\x20-\x7E\n\r\t]', '', content)
                    content = re.sub(r'\s+', ' ', content)
            except Exception as e2:
                raise Exception(f"Could not parse DOC file: {e2}")
        
        return {
            'content': content.strip(),
            'metadata': metadata,
            'total_words': len(content.split()),
            'parsing_method': 'external_tool_fallback'
        }
    
    def _parse_text(self, file_path: Path) -> Dict[str, Any]:
        """Parse plain text or markdown file.
        
        Args:
            file_path: Path to text file
            
        Returns:
            Dictionary with extracted content
        """
        with open(file_path, 'r', encoding='utf-8') as f:
            content = f.read()
        
        return {
            'content': content,
            'metadata': {},
            'total_words': len(content.split()),
            'parsing_method': 'direct_read'
        }
    
    def _parse_rtf(self, file_path: Path) -> Dict[str, Any]:
        """Parse RTF document.
        
        Args:
            file_path: Path to RTF file
            
        Returns:
            Dictionary with extracted content
        """
        try:
            import striprtf
        except ImportError:
            raise ImportError(
                "striprtf is required for RTF parsing. Install with: pip install striprtf"
            )
        
        with open(file_path, 'r', encoding='utf-8') as f:
            rtf_content = f.read()
        
        # Strip RTF formatting
        content = striprtf.striprtf(rtf_content)
        
        return {
            'content': content.strip(),
            'metadata': {},
            'total_words': len(content.split()),
            'parsing_method': 'striprtf'
        }
    
    def extract_citations_from_document(self, file_path: Union[str, Path]) -> List[Dict[str, Any]]:
        """Extract citations from document with page/paragraph information.
        
        Args:
            file_path: Path to document file
            
        Returns:
            List of citations with location information
        """
        parsed_doc = self.parse_document(file_path)
        content = parsed_doc['content']
        
        # Import citation extraction from NLP processor
        from .text_processing import extract_citations
        
        citations = extract_citations(content)
        citations_with_location = []
        
        for citation in citations:
            # Find citation position in content
            citation_pos = content.find(citation)
            
            citation_info = {
                'text': citation,
                'position': citation_pos,
                'context': self._extract_context_around_position(content, citation_pos)
            }
            
            # Add page/paragraph information if available
            if 'pages' in parsed_doc:
                citation_info['page'] = self._find_page_for_position(citation_pos, parsed_doc['pages'])
            
            if 'paragraphs' in parsed_doc:
                citation_info['paragraph'] = self._find_paragraph_for_position(citation_pos, parsed_doc['paragraphs'])
            
            citations_with_location.append(citation_info)
        
        return citations_with_location
    
    def _extract_context_around_position(self, content: str, position: int, window: int = 200) -> str:
        """Extract context around a specific position.
        
        Args:
            content: Full content
            position: Position of the citation
            window: Number of characters before and after
            
        Returns:
            Context string
        """
        if position == -1:
            return ""
        
        start = max(0, position - window)
        end = min(len(content), position + window)
        
        return content[start:end].strip()
    
    def _find_page_for_position(self, position: int, pages: List[Dict[str, Any]]) -> Optional[int]:
        """Find which page contains the given position.
        
        Args:
            position: Character position in document
            pages: List of page information
            
        Returns:
            Page number or None
        """
        current_pos = 0
        
        for page in pages:
            page_length = len(page['text'])
            if current_pos <= position < current_pos + page_length:
                return page['page_number']
            current_pos += page_length
        
        return None
    
    def _find_paragraph_for_position(self, position: int, paragraphs: List[Dict[str, Any]]) -> Optional[int]:
        """Find which paragraph contains the given position.
        
        Args:
            position: Character position in document
            paragraphs: List of paragraph information
            
        Returns:
            Paragraph number or None
        """
        current_pos = 0
        
        for paragraph in paragraphs:
            paragraph_length = len(paragraph['text'])
            if current_pos <= position < current_pos + paragraph_length:
                return paragraph['paragraph_number']
            current_pos += paragraph_length
        
        return None
    
    def _parse_pptx(self, file_path: Path) -> Dict[str, Any]:
        """Parse PowerPoint PPTX document.
        
        Args:
            file_path: Path to PPTX file
            
        Returns:
            Dictionary with extracted content
        """
        try:
            from pptx import Presentation
        except ImportError:
            raise ImportError(
                "python-pptx is required for PPTX parsing. Install with: pip install python-pptx"
            )
        
        prs = Presentation(file_path)
        content = ""
        slides = []
        
        # Extract metadata
        metadata = {}
        if hasattr(prs.core_properties, 'title'):
            metadata = {
                'title': prs.core_properties.title or '',
                'author': prs.core_properties.author or '',
                'subject': prs.core_properties.subject or '',
                'keywords': prs.core_properties.keywords or '',
                'created': str(prs.core_properties.created) if prs.core_properties.created else '',
                'modified': str(prs.core_properties.modified) if prs.core_properties.modified else ''
            }
        
        # Extract content from slides
        for slide_num, slide in enumerate(prs.slides):
            slide_content = ""
            slide_texts = []
            
            # Extract text from all shapes in the slide
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text = shape.text.strip()
                    slide_texts.append(text)
                    slide_content += text + "\n"
            
            slides.append({
                'slide_number': slide_num + 1,
                'text': slide_content.strip(),
                'text_blocks': slide_texts,
                'word_count': len(slide_content.split())
            })
            content += slide_content + "\n\n"
        
        return {
            'content': content.strip(),
            'metadata': metadata,
            'slides': slides,
            'total_slides': len(slides),
            'total_words': len(content.split()),
            'parsing_method': 'python-pptx'
        }
    
    def _parse_ppt(self, file_path: Path) -> Dict[str, Any]:
        """Parse PowerPoint PPT document (legacy format).
        
        Args:
            file_path: Path to PPT file
            
        Returns:
            Dictionary with extracted content
        """
        # For legacy PPT files, we'll try to extract what we can
        # This is more limited than PPTX parsing
        content = ""
        metadata = {}
        
        try:
            # Try to read as binary and extract text patterns
            with open(file_path, 'rb') as f:
                raw_content = f.read()
                
            # Convert to string and clean up
            text_content = raw_content.decode('utf-8', errors='ignore')
            
            # Extract readable text (very basic approach)
            import re
            # Remove binary data and keep only readable characters
            text_content = re.sub(r'[^\x20-\x7E\n\r\t]', ' ', text_content)
            # Clean up multiple spaces
            text_content = re.sub(r'\s+', ' ', text_content)
            
            # Look for title-like patterns
            title_pattern = r'([A-Z][A-Z\s]{10,})'
            potential_titles = re.findall(title_pattern, text_content)
            
            content = text_content.strip()
            if potential_titles:
                metadata['potential_titles'] = potential_titles[:5]  # First 5 potential titles
            
        except Exception as e:
            logger.warning(f"Could not parse PPT file {file_path}: {e}")
            content = f"[Error parsing PPT file: {e}]"
        
        return {
            'content': content,
            'metadata': metadata,
            'total_words': len(content.split()),
            'parsing_method': 'basic_text_extraction',
            'note': 'Legacy PPT format - limited text extraction'
        }
    
    def _parse_media(self, file_path: Path) -> Dict[str, Any]:
        """Parse media files (audio/video) to extract text/transcripts.
        
        Args:
            file_path: Path to media file
            
        Returns:
            Dictionary with extracted content or metadata
        """
        file_extension = file_path.suffix.lower()
        file_size = file_path.stat().st_size
        content = ""
        metadata = {
            'file_type': 'media',
            'file_extension': file_extension,
            'file_size': file_size,
            'parsing_method': 'metadata_only'
        }
        
        # Check if file is very large and warn user
        if file_size > 200 * 1024 * 1024:  # 200MB
            size_mb = file_size / (1024 * 1024)
            logger.warning(f"Processing large media file: {file_path.name} ({size_mb:.1f} MB). This may take a while...")
        
        # For audio files, try to extract transcript
        if file_extension in ['.mp3', '.wav', '.m4a']:
            content = self._extract_audio_transcript(file_path)
            metadata['parsing_method'] = 'speech_recognition'
        
        # For video files, try to extract audio and then transcript
        elif file_extension in ['.mp4', '.mov', '.avi']:
            content = self._extract_video_transcript(file_path)
            metadata['parsing_method'] = 'video_audio_extraction'
        
        # If we couldn't extract text, provide file metadata
        if not content.strip():
            content = f"[Media file: {file_path.name}]\n"
            content += f"File type: {file_extension}\n"
            content += f"Size: {file_size / (1024*1024):.1f} MB\n"
            content += f"Note: Text extraction not available or failed for this media format."
        
        return {
            'content': content.strip(),
            'metadata': metadata,
            'total_words': len(content.split()),
            'parsing_method': metadata['parsing_method']
        }
    
    def _extract_audio_transcript(self, file_path: Path) -> str:
        """Extract transcript from audio file using speech recognition.
        
        Args:
            file_path: Path to audio file
            
        Returns:
            Extracted transcript text
        """
        try:
            import speech_recognition as sr
            from pydub import AudioSegment
        except ImportError:
            return f"[Audio transcript extraction not available - install SpeechRecognition and pydub]"
        
        try:
            # Initialize speech recognizer
            recognizer = sr.Recognizer()
            
            # Convert audio to WAV format if needed
            audio_file = str(file_path)
            if file_path.suffix.lower() != '.wav':
                # Convert to WAV using pydub
                # pydub can handle m4a files - it uses ffmpeg under the hood
                try:
                    # Try to load with explicit format for m4a files
                    if file_path.suffix.lower() == '.m4a':
                        audio = AudioSegment.from_file(file_path, format="m4a")
                    else:
                        audio = AudioSegment.from_file(file_path)
                    
                    # Create temporary WAV file
                    import tempfile
                    import os
                    temp_wav = tempfile.NamedTemporaryFile(suffix='.wav', delete=False)
                    temp_wav.close()
                    audio.export(temp_wav.name, format="wav")
                    audio_file = temp_wav.name
                except Exception as e:
                    logger.warning(f"Could not convert {file_path.suffix} to WAV: {e}")
                    # Fallback: try without format specification
                    try:
                        audio = AudioSegment.from_file(str(file_path))
                        import tempfile
                        import os
                        temp_wav = tempfile.NamedTemporaryFile(suffix='.wav', delete=False)
                        temp_wav.close()
                        audio.export(temp_wav.name, format="wav")
                        audio_file = temp_wav.name
                    except Exception as e2:
                        logger.error(f"Failed to convert audio file {file_path}: {e2}")
                        return f"[Audio conversion failed: {e2}]"
            
            # Extract transcript
            try:
                with sr.AudioFile(audio_file) as source:
                    audio_data = recognizer.record(source)
                    transcript = recognizer.recognize_google(audio_data)
                    return transcript
            finally:
                # Clean up temporary WAV file if we created one
                if audio_file != str(file_path) and os.path.exists(audio_file):
                    try:
                        os.unlink(audio_file)
                    except Exception as e:
                        logger.warning(f"Could not delete temporary audio file {audio_file}: {e}")
                
        except Exception as e:
            logger.warning(f"Could not extract audio transcript from {file_path}: {e}")
            # Clean up temp file on error
            if 'temp_wav' in locals() and 'audio_file' in locals() and audio_file != str(file_path):
                try:
                    if os.path.exists(audio_file):
                        os.unlink(audio_file)
                except:
                    pass
            return f"[Audio transcript extraction failed: {e}]"
    
    def _extract_video_transcript(self, file_path: Path) -> str:
        """Extract transcript from video file by extracting audio first.
        
        Args:
            file_path: Path to video file
            
        Returns:
            Extracted transcript text
        """
        try:
            from moviepy.editor import VideoFileClip
            import tempfile
        except ImportError:
            return f"[Video transcript extraction not available - install moviepy]"
        
        try:
            # Extract audio from video
            video = VideoFileClip(str(file_path))
            audio = video.audio
            
            # Save audio to temporary file
            with tempfile.NamedTemporaryFile(suffix='.wav', delete=False) as temp_audio:
                audio.write_audiofile(temp_audio.name, verbose=False, logger=None)
                temp_audio_path = Path(temp_audio.name)
            
            # Extract transcript from audio
            transcript = self._extract_audio_transcript(temp_audio_path)
            
            # Clean up
            video.close()
            audio.close()
            temp_audio_path.unlink(missing_ok=True)
            
            return transcript
            
        except Exception as e:
            logger.warning(f"Could not extract video transcript from {file_path}: {e}")
            return f"[Video transcript extraction failed: {e}]"
    
    def get_supported_formats(self) -> List[str]:
        """Get list of supported file formats.
        
        Returns:
            List of supported file extensions
        """
        return list(self.supported_formats.keys())
    
    def is_supported(self, file_path: Union[str, Path]) -> bool:
        """Check if file format is supported.
        
        Args:
            file_path: Path to file
            
        Returns:
            True if format is supported
        """
        path = Path(file_path)
        return path.suffix.lower() in self.supported_formats
